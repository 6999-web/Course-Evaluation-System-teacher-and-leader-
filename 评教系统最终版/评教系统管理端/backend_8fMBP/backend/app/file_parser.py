"""
文件解析器
用于提取不同格式文件的文本内容
"""

import logging
import os
from typing import Optional

logger = logging.getLogger(__name__)


class FileParser:
    """文件解析器"""
    
    @staticmethod
    def parse_docx(file_path: str) -> str:
        """
        解析 DOCX 文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            提取的文本内容
            
        Raises:
            Exception: 解析失败
        """
        try:
            from docx import Document
            
            doc = Document(file_path)
            text_content = []
            
            # 提取段落文本
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # 提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text)
                    if row_text:
                        text_content.append(" | ".join(row_text))
            
            return "\n".join(text_content)
            
        except ImportError:
            raise Exception("缺少 python-docx 库，请安装: pip install python-docx")
        except Exception as e:
            # 如果DOCX解析失败，尝试作为纯文本文件读取
            # 有些文件虽然扩展名是.docx，但实际上是纯文本文件
            logger.warning(f"DOCX解析失败，尝试作为纯文本读取: {str(e)}")
            try:
                return FileParser.parse_txt(file_path)
            except Exception as txt_error:
                raise Exception(f"DOCX 文件解析失败: {str(e)}，纯文本读取也失败: {str(txt_error)}")
    
    @staticmethod
    def parse_pdf(file_path: str) -> str:
        """
        解析 PDF 文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            提取的文本内容
            
        Raises:
            Exception: 解析失败
        """
        try:
            import PyPDF2
            
            text_content = []
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text.strip():
                        text_content.append(text)
            
            return "\n".join(text_content)
            
        except ImportError:
            raise Exception("缺少 PyPDF2 库，请安装: pip install PyPDF2")
        except Exception as e:
            raise Exception(f"PDF 文件解析失败: {str(e)}")
    
    @staticmethod
    def parse_pptx(file_path: str) -> str:
        """
        解析 PPTX 文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            提取的文本内容
            
        Raises:
            Exception: 解析失败
        """
        try:
            from pptx import Presentation
            
            text_content = []
            
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                
                # 提取文本框内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_content.append(f"[幻灯片 {slide_num + 1}]\n" + "\n".join(slide_text))
            
            return "\n".join(text_content)
            
        except ImportError:
            raise Exception("缺少 python-pptx 库，请安装: pip install python-pptx")
        except Exception as e:
            raise Exception(f"PPTX 文件解析失败: {str(e)}")
    
    @staticmethod
    def parse_txt(file_path: str) -> str:
        """
        解析 TXT 文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            提取的文本内容
            
        Raises:
            Exception: 解析失败
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                with open(file_path, 'r', encoding='gbk') as file:
                    return file.read()
            except Exception as e:
                raise Exception(f"TXT 文件解析失败: {str(e)}")
        except Exception as e:
            raise Exception(f"TXT 文件解析失败: {str(e)}")
    
    @staticmethod
    def parse_file(file_path: str, file_type: Optional[str] = None) -> str:
        """
        根据文件类型自动选择解析方法
        
        Args:
            file_path: 文件路径
            file_type: 文件类型（可选，如果不提供则从文件扩展名推断）
            
        Returns:
            提取的文本内容
            
        Raises:
            Exception: 不支持的文件格式或解析失败
        """
        try:
            # 检查文件是否存在
            if not os.path.exists(file_path):
                raise Exception(f"文件不存在: {file_path}")
            
            # 检查文件大小
            file_size = os.path.getsize(file_path)
            if file_size == 0:
                raise Exception("文件内容为空")
            
            # 如果没有指定文件类型，从扩展名推断
            if not file_type:
                _, ext = os.path.splitext(file_path)
                file_type = ext.lower().lstrip('.')
            
            # 根据文件类型选择解析方法
            if file_type in ['docx', 'doc']:
                return FileParser.parse_docx(file_path)
            elif file_type == 'pdf':
                return FileParser.parse_pdf(file_path)
            elif file_type in ['pptx', 'ppt']:
                return FileParser.parse_pptx(file_path)
            elif file_type == 'txt':
                return FileParser.parse_txt(file_path)
            else:
                raise Exception(f"不支持的文件格式: {file_type}。支持的格式: docx, pdf, pptx, txt")
            
        except Exception as e:
            logger.error(f"文件解析失败: {str(e)}")
            raise

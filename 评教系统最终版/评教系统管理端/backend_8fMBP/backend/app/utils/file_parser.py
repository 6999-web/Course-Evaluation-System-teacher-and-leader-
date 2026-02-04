"""
文件解析器模块

支持解析以下格式的文件：
- DOCX (Word文档)
- PDF (PDF文档)
- PPTX (PowerPoint演示文稿)

提取文件中的文本内容用于自动评分
"""

import os
import logging
from typing import Optional, Dict
from pathlib import Path

# 文件解析库
try:
    from docx import Document
except ImportError:
    Document = None

try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# 配置日志
logger = logging.getLogger(__name__)


class FileParserError(Exception):
    """文件解析错误基类"""
    pass


class UnsupportedFormatError(FileParserError):
    """不支持的文件格式错误"""
    pass


class EmptyFileError(FileParserError):
    """空文件错误"""
    pass


class FileCorruptedError(FileParserError):
    """文件损坏错误"""
    pass


class FileParser:
    """
    文件解析器类
    
    负责提取不同格式文件的文本内容
    """
    
    # 支持的文件格式
    SUPPORTED_FORMATS = {
        'docx': 'Word文档',
        'pdf': 'PDF文档',
        'pptx': 'PowerPoint演示文稿',
        'ppt': 'PowerPoint演示文稿'
    }
    
    def __init__(self):
        """初始化文件解析器"""
        self._check_dependencies()
    
    def _check_dependencies(self):
        """检查必需的依赖库是否已安装"""
        missing_deps = []
        
        if Document is None:
            missing_deps.append("python-docx")
        if PdfReader is None:
            missing_deps.append("PyPDF2")
        if Presentation is None:
            missing_deps.append("python-pptx")
        
        if missing_deps:
            logger.warning(f"缺少依赖库: {', '.join(missing_deps)}")
    
    def parse_file(self, file_path: str, file_type: Optional[str] = None) -> str:
        """
        根据文件类型自动选择解析方法
        
        Args:
            file_path: 文件路径
            file_type: 文件类型（可选，如果不提供则从文件扩展名推断）
        
        Returns:
            str: 提取的文本内容
        
        Raises:
            UnsupportedFormatError: 不支持的文件格式
            EmptyFileError: 文件内容为空
            FileCorruptedError: 文件损坏
            FileNotFoundError: 文件不存在
        """
        # 检查文件是否存在
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        if file_size == 0:
            raise EmptyFileError(f"文件为空: {file_path}")
        
        # 确定文件类型
        if file_type is None:
            file_type = Path(file_path).suffix.lower().lstrip('.')
        else:
            file_type = file_type.lower()
        
        # 检查是否支持该格式
        if file_type not in self.SUPPORTED_FORMATS:
            raise UnsupportedFormatError(
                f"不支持的文件格式: {file_type}。"
                f"支持的格式: {', '.join(self.SUPPORTED_FORMATS.keys())}"
            )
        
        # 根据文件类型选择解析方法
        try:
            if file_type == 'docx':
                content = self.parse_docx(file_path)
            elif file_type == 'pdf':
                content = self.parse_pdf(file_path)
            elif file_type in ['pptx', 'ppt']:
                content = self.parse_pptx(file_path)
            else:
                raise UnsupportedFormatError(f"不支持的文件格式: {file_type}")
            
            # 检查提取的内容是否为空
            if not content or content.strip() == "":
                raise EmptyFileError(f"文件内容为空或无法提取文本: {file_path}")
            
            logger.info(f"成功解析文件: {file_path}, 提取文本长度: {len(content)}")
            return content
            
        except (UnsupportedFormatError, EmptyFileError):
            raise
        except Exception as e:
            logger.error(f"解析文件失败: {file_path}, 错误: {str(e)}")
            raise FileCorruptedError(f"文件可能已损坏或格式不正确: {str(e)}")
    
    def parse_docx(self, file_path: str) -> str:
        """
        解析 DOCX 文件
        
        Args:
            file_path: DOCX 文件路径
        
        Returns:
            str: 提取的文本内容
        
        Raises:
            FileCorruptedError: 文件损坏或无法解析
        """
        if Document is None:
            raise UnsupportedFormatError("未安装 python-docx 库，无法解析 DOCX 文件")
        
        try:
            doc = Document(file_path)
            
            # 提取所有段落的文本
            paragraphs = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)
            
            # 提取表格中的文本
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            paragraphs.append(text)
            
            content = '\n'.join(paragraphs)
            
            logger.debug(f"DOCX解析完成: {file_path}, 段落数: {len(paragraphs)}")
            return content
            
        except Exception as e:
            logger.error(f"DOCX解析失败: {file_path}, 错误: {str(e)}")
            raise FileCorruptedError(f"DOCX文件解析失败: {str(e)}")
    
    def parse_pdf(self, file_path: str) -> str:
        """
        解析 PDF 文件
        
        Args:
            file_path: PDF 文件路径
        
        Returns:
            str: 提取的文本内容
        
        Raises:
            FileCorruptedError: 文件损坏或无法解析
        """
        if PdfReader is None:
            raise UnsupportedFormatError("未安装 PyPDF2 库，无法解析 PDF 文件")
        
        try:
            reader = PdfReader(file_path)
            
            # 提取所有页面的文本
            pages_text = []
            for page_num, page in enumerate(reader.pages):
                try:
                    text = page.extract_text()
                    if text and text.strip():
                        pages_text.append(text.strip())
                except Exception as e:
                    logger.warning(f"PDF第{page_num+1}页提取失败: {str(e)}")
                    continue
            
            content = '\n'.join(pages_text)
            
            logger.debug(f"PDF解析完成: {file_path}, 页数: {len(reader.pages)}, 提取页数: {len(pages_text)}")
            return content
            
        except Exception as e:
            logger.error(f"PDF解析失败: {file_path}, 错误: {str(e)}")
            raise FileCorruptedError(f"PDF文件解析失败: {str(e)}")
    
    def parse_pptx(self, file_path: str) -> str:
        """
        解析 PPTX 文件
        
        Args:
            file_path: PPTX 文件路径
        
        Returns:
            str: 提取的文本内容
        
        Raises:
            FileCorruptedError: 文件损坏或无法解析
        """
        if Presentation is None:
            raise UnsupportedFormatError("未安装 python-pptx 库，无法解析 PPTX 文件")
        
        try:
            prs = Presentation(file_path)
            
            # 提取所有幻灯片的文本
            slides_text = []
            for slide_num, slide in enumerate(prs.slides):
                slide_content = []
                
                # 提取幻灯片中所有形状的文本
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_content.append(text)
                    
                    # 如果是表格，提取表格内容
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                text = cell.text.strip()
                                if text:
                                    slide_content.append(text)
                
                if slide_content:
                    slides_text.append(f"[幻灯片 {slide_num + 1}]\n" + '\n'.join(slide_content))
            
            content = '\n\n'.join(slides_text)
            
            logger.debug(f"PPTX解析完成: {file_path}, 幻灯片数: {len(prs.slides)}")
            return content
            
        except Exception as e:
            logger.error(f"PPTX解析失败: {file_path}, 错误: {str(e)}")
            raise FileCorruptedError(f"PPTX文件解析失败: {str(e)}")
    
    def get_file_info(self, file_path: str) -> Dict[str, any]:
        """
        获取文件基本信息
        
        Args:
            file_path: 文件路径
        
        Returns:
            dict: 文件信息字典
        """
        import datetime
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        file_stat = os.stat(file_path)
        file_ext = Path(file_path).suffix.lower().lstrip('.')
        
        # 将时间戳转换为字符串格式
        modified_time = datetime.datetime.fromtimestamp(file_stat.st_mtime).isoformat()
        
        return {
            'file_name': os.path.basename(file_path),
            'file_path': file_path,
            'file_type': file_ext,
            'file_size': file_stat.st_size,
            'file_size_mb': round(file_stat.st_size / (1024 * 1024), 2),
            'is_supported': file_ext in self.SUPPORTED_FORMATS,
            'format_name': self.SUPPORTED_FORMATS.get(file_ext, '未知格式'),
            'modified_time': modified_time
        }
    
    @staticmethod
    def is_supported_format(file_type: str) -> bool:
        """
        检查文件格式是否支持
        
        Args:
            file_type: 文件类型（扩展名）
        
        Returns:
            bool: 是否支持
        """
        return file_type.lower() in FileParser.SUPPORTED_FORMATS


# 创建全局解析器实例
file_parser = FileParser()


# 便捷函数
def parse_file(file_path: str, file_type: Optional[str] = None) -> str:
    """
    解析文件的便捷函数
    
    Args:
        file_path: 文件路径
        file_type: 文件类型（可选）
    
    Returns:
        str: 提取的文本内容
    """
    return file_parser.parse_file(file_path, file_type)


def get_file_info(file_path: str) -> Dict[str, any]:
    """
    获取文件信息的便捷函数
    
    Args:
        file_path: 文件路径
    
    Returns:
        dict: 文件信息
    """
    return file_parser.get_file_info(file_path)


def is_supported_format(file_type: str) -> bool:
    """
    检查文件格式是否支持的便捷函数
    
    Args:
        file_type: 文件类型
    
    Returns:
        bool: 是否支持
    """
    return FileParser.is_supported_format(file_type)
